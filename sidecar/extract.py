"""
非 PDF 文档的文本提取(Office / Markdown / 纯文本)。
统一返回 [(unit_no, text, method), ...],unit 相当于 PDF 的"页":
  - docx: 按 ~1500 字切块
  - pptx: 每张幻灯片一个 unit
  - xlsx: 每个工作表一个 unit
  - md/txt: 按 ~1500 字切块
PDF/EPUB 等仍走 fitz(见 ingest.process_pdf)。
"""
from __future__ import annotations
import os

CHUNK = 1500


def _chunk_text(text: str, size: int = CHUNK):
    text = (text or "").strip()
    if not text:
        return []
    return [text[i:i + size] for i in range(0, len(text), size)]


def _docx(path):
    from docx import Document
    doc = Document(path)
    paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    # 表格文字也带上
    for t in doc.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                paras.append(" | ".join(cells))
    full = "\n".join(paras)
    return [(i + 1, chunk, "docx") for i, chunk in enumerate(_chunk_text(full))]


def _pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    units = []
    for i, slide in enumerate(prs.slides):
        buf = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in para.runs) or para.text
                    if line and line.strip():
                        buf.append(line.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                    if cells:
                        buf.append(" | ".join(cells))
        txt = "\n".join(buf).strip()
        units.append((i + 1, txt or "(空白幻灯片)", "pptx"))
    return units


def _xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    units = []
    for i, ws in enumerate(wb.worksheets):
        rows = []
        for row in ws.iter_rows(values_only=True):
            vals = [str(v) for v in row if v is not None]
            if vals:
                rows.append(" | ".join(vals))
        txt = f"# 工作表: {ws.title}\n" + "\n".join(rows)
        units.append((i + 1, txt.strip(), "xlsx"))
    wb.close()
    return units


def _text(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    ext = os.path.splitext(path)[1].lower()
    method = "markdown" if ext in (".md", ".markdown") else "text"
    return [(i + 1, chunk, method) for i, chunk in enumerate(_chunk_text(raw))]


def _html(path):
    """网页 / 聊天记录 html 导出 → 去标签取正文。"""
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for s in soup(["script", "style", "noscript"]):
        s.decompose()
    lines = [ln.strip() for ln in soup.get_text("\n").splitlines() if ln.strip()]
    return [(i + 1, chunk, "html") for i, chunk in enumerate(_chunk_text("\n".join(lines)))]


def _csv(path):
    import csv
    rows = []
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        for row in csv.reader(f):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
    return [(i + 1, chunk, "csv") for i, chunk in enumerate(_chunk_text("\n".join(rows)))]


def _json(path):
    """聊天记录 / 结构化 json → 递归取所有字符串值。"""
    import json
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        data = json.load(f)
    out = []

    def walk(x):
        if isinstance(x, str):
            out.append(x)
        elif isinstance(x, dict):
            for v in x.values():
                walk(v)
        elif isinstance(x, list):
            for v in x:
                walk(v)
    walk(data)
    full = "\n".join(s.strip() for s in out if s and s.strip())
    return [(i + 1, chunk, "json") for i, chunk in enumerate(_chunk_text(full))]


def _safe_body(part):
    """取邮件正文,治中文乱码:邮件未声明 charset 时 get_content() 按 us-ascii 解码 → UTF-8 中文变 �。
    检测到无 charset 或出现替换字符,就用原始字节按 utf-8/gbk/big5 兜底重解。"""
    if part is None:
        return ""
    txt = None
    try:
        txt = part.get_content()
    except Exception:
        txt = None
    try:
        charset = part.get_content_charset()
    except Exception:
        charset = None
    if txt is None or "�" in txt or not charset:
        try:
            raw = part.get_payload(decode=True)   # 先解 base64/quoted-printable 回原始字节
            if raw:
                for enc in ("utf-8", "gbk", "gb18030", "big5"):
                    try:
                        d = raw.decode(enc)
                        if "�" not in d:
                            return d
                    except Exception:
                        pass
                return raw.decode("utf-8", errors="replace")
        except Exception:
            pass
    return txt or ""


def _mbox(path):
    """Gmail Takeout 导出的 .mbox:一个文件含全部邮件。每封一个 unit。"""
    import mailbox
    import email
    import email.policy

    def factory(f):
        return email.message_from_binary_file(f, policy=email.policy.default)

    units = []
    mb = mailbox.mbox(path, factory=factory)
    for i, msg in enumerate(mb):
        body = ""
        try:
            part = msg.get_body(preferencelist=("plain", "html"))
            body = _safe_body(part)
        except Exception:
            body = ""
        txt = f"主题: {msg.get('subject', '')}\n发件: {msg.get('from', '')}\n日期: {msg.get('date', '')}\n\n{body}"
        units.append((i + 1, txt.strip()[:8000], "email"))
    return units if units else [(1, "(空邮箱)", "email")]


def _eml(path):
    from email import policy
    from email.parser import BytesParser
    with open(path, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)
    body = ""
    try:
        part = msg.get_body(preferencelist=("plain", "html"))
        body = _safe_body(part) or _safe_body(msg)
    except Exception:
        body = _safe_body(msg)
    full = f"主题: {msg.get('subject', '')}\n发件: {msg.get('from', '')}\n\n{body}"
    return [(i + 1, chunk, "email") for i, chunk in enumerate(_chunk_text(full))]



# ===== 微信聊天记录解析(MemoTrace/留痕 等工具导出的 CSV/TXT) =====
# 消息类型映射(微信 Type 字段):还原成可读占位,而非丢弃
_WX_TYPE = {
    "1": "", "3": "[图片]", "34": "[语音]", "43": "[视频]", "42": "[名片]",
    "47": "[表情]", "48": "[位置]", "49": "[链接/文件]", "50": "[语音/视频通话]",
    "10000": "[系统消息]", "436207665": "[红包]", "419430449": "[转账]",
}
_WX_COLS = ("strcontent", "iscender", "issender", "createtime", "strtime", "nickname", "type")


def _wx_contact(lines):
    """从对话行里认出"对话对象"(出现最多的、非"我"的发言人)。"""
    import re as _re
    from collections import Counter
    c = Counter()
    for ln in lines:
        # 行形如 "[时间] 名: 内容" 或 "名: 内容" → 跳过[时间]再取名
        m = _re.match(r"^(?:\[[^\]]*\]\s*)?([^:\n]{1,20}?)\s*:", ln)
        if m:
            nm = m.group(1).strip()
            if nm and nm not in ("我", "对方"):
                c[nm] += 1
    return c.most_common(1)[0][0] if c else "对方"


def _wx_finalize(lines):
    """认人:识别对话对象 → 加标题 + 每个分块打「与X」标签,让"和某某的聊天"能被检索锁定。"""
    if not lines:
        return [(1, "(空的微信记录)", "wechat")]
    contact = _wx_contact(lines)
    header = "微信聊天记录 · 与「%s」的对话(共 %d 条消息)" % (contact, len(lines))
    tag = "【微信 · 与%s的聊天】" % contact
    full = header + "\n" + "\n".join(lines)
    return [(i + 1, (chunk if i == 0 else tag + "\n" + chunk), "wechat")
            for i, chunk in enumerate(_chunk_text(full))]


def _looks_wechat_csv(path):
    """看 CSV 头是否是微信导出(有 StrContent + IsSender + 时间列)。"""
    import csv
    try:
        with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
            head = next(csv.reader(f))
        low = [h.strip().lower() for h in head]
        return ("strcontent" in low) and ("issender" in low) and any(k in low for k in ("createtime", "strtime"))
    except Exception:
        return False


def _wechat_csv(path):
    """微信 CSV(MemoTrace):还原成 [时间] 发言人: 内容 的对话流。"""
    import csv
    lines = []
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        rd = csv.DictReader(f)
        low = {k.lower(): k for k in (rd.fieldnames or [])}
        c_content = low.get("strcontent")
        c_is = low.get("issender")
        c_time = low.get("strtime") or low.get("createtime")
        c_name = low.get("nickname") or low.get("remark") or low.get("sender") or low.get("talker")
        c_type = low.get("type")
        for row in rd:
            content = (row.get(c_content) or "").strip() if c_content else ""
            mtype = (row.get(c_type) or "1").strip() if c_type else "1"
            placeholder = _WX_TYPE.get(mtype, "")
            if placeholder and not content:
                content = placeholder
            elif placeholder and content:
                content = placeholder + " " + content
            if not content:
                continue
            is_me = str(row.get(c_is) or "0").strip() in ("1", "1.0", "true", "True") if c_is else False
            name = "我" if is_me else ((row.get(c_name) or "对方").strip() if c_name else "对方")
            t = (row.get(c_time) or "").strip() if c_time else ""
            lines.append(("[%s] %s: %s" % (t, name, content)) if t else ("%s: %s" % (name, content)))
    return _wx_finalize(lines)


def _looks_wechat_txt(path):
    """MemoTrace TXT 导出:每条含 '昵称  时间' 行。取样前几十行判断。"""
    import re as _re
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            head = f.read(4000)
        # 形如 "某某  2024-01-01 12:00:00" 的行出现多次
        hits = _re.findall(r"^.{1,30}\s+\d{4}-\d{1,2}-\d{1,2}\s+\d{1,2}:\d{2}", head, _re.M)
        return len(hits) >= 2
    except Exception:
        return False


def _wechat_txt(path):
    """MemoTrace TXT:'昵称  时间\n内容' 块 → 规整成对话流。"""
    import re as _re
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    blocks = []
    cur_head = None
    cur_body = []
    for ln in raw.splitlines():
        m = _re.match(r"^(.{1,30}?)\s+(\d{4}-\d{1,2}-\d{1,2}\s+\d{1,2}:\d{2}(?::\d{2})?)\s*$", ln)
        if m:
            if cur_head and cur_body:
                blocks.append("[%s] %s: %s" % (cur_head[1], cur_head[0], " ".join(cur_body).strip()))
            cur_head = (m.group(1).strip(), m.group(2).strip()); cur_body = []
        elif ln.strip():
            cur_body.append(ln.strip())
    if cur_head and cur_body:
        blocks.append("[%s] %s: %s" % (cur_head[1], cur_head[0], " ".join(cur_body).strip()))
    if blocks:
        return _wx_finalize(blocks)
    return [(i + 1, chunk, "wechat") for i, chunk in enumerate(_chunk_text(raw))]


def extract_units(path: str):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        return _docx(path)
    if ext == ".pptx":
        return _pptx(path)
    if ext in (".xlsx", ".xlsm"):
        return _xlsx(path)
    if ext == ".txt" and _looks_wechat_txt(path):
        return _wechat_txt(path)
    if ext in (".md", ".markdown", ".txt"):
        return _text(path)
    if ext in (".html", ".htm"):
        return _html(path)
    if ext == ".csv":
        if _looks_wechat_csv(path):
            return _wechat_csv(path)
        return _csv(path)
    if ext == ".json":
        return _json(path)
    if ext == ".eml":
        return _eml(path)
    if ext == ".mbox":
        return _mbox(path)
    return []

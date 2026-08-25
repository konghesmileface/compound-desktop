"""
产出机:基于用户知识库检索的素材 → LLM 组织结构 → 渲染成真文件(PPT/Word/Excel)。
格式无关:同一套"检索+组织",不同渲染器输出。PPT 支持多套设计主题(模板可选)。
"""
from __future__ import annotations
import os
import re
import json

BRAIN = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(BRAIN, "generated")
os.makedirs(OUT_DIR, exist_ok=True)

# ===== PPT 设计主题(RGB)。前端可选;default=deep =====
THEMES = {
    "deep":  {"name": "深空", "bg": (14, 17, 26), "title": (232, 236, 243), "accent": (103, 232, 249), "body": (183, 190, 202), "foot": (90, 98, 112), "font": "Microsoft YaHei"},
    "clean": {"name": "简约", "bg": (255, 255, 255), "title": (24, 28, 38), "accent": (37, 99, 235), "body": (60, 66, 78), "foot": (150, 156, 166), "font": "Microsoft YaHei"},
    "warm":  {"name": "暖阳", "bg": (252, 248, 240), "title": (60, 40, 22), "accent": (234, 128, 42), "body": (86, 70, 54), "foot": (172, 158, 140), "font": "Microsoft YaHei"},
    "forest":{"name": "松林", "bg": (18, 30, 26), "title": (232, 244, 238), "accent": (52, 211, 153), "body": (176, 196, 188), "foot": (96, 116, 108), "font": "Microsoft YaHei"},
}
DEFAULT_THEME = "deep"


def _parse_json(text: str):
    """从 LLM 回复里抠出 JSON(容忍 markdown 代码围栏)。"""
    m = re.search(r"```(?:json)?\s*(\{.*\})\s*```", text, re.S)
    raw = m.group(1) if m else text
    s, e = raw.find("{"), raw.rfind("}")
    if s >= 0 and e > s:
        raw = raw[s:e + 1]
    return json.loads(raw)


def compose(LLM, topic: str, sources: list, fmt: str) -> dict:
    """让 LLM 基于 sources 为 topic 组织结构化内容(JSON)。"""
    context = "\n\n".join(
        f"[{i + 1}] {s['filename']} 第{s['page_no']}页:\n{s['text']}"
        for i, s in enumerate(sources))
    if fmt == "excel":
        schema = '{"title":"表标题","columns":["列1","列2",...],"rows":[["...","..."],...]}'
        what = "一份数据/清单表格,列名精确、每行是一条真实可用的记录"
    elif fmt == "word":
        schema = ('{"title":"文档标题","subtitle":"一句话副标题",'
                  '"sections":[{"heading":"小节标题","paragraphs":["段落1","段落2"]},...]}')
        what = "一份结构化文档(4-6 个小节,每节 2-3 段,论述充分)"
    else:  # ppt
        schema = ('{"title":"演示标题","subtitle":"一句话副标题",'
                  '"slides":[{"title":"页标题","bullets":["要点1","要点2","要点3"],"note":"这一页的核心结论(一句话)"},...]}')
        what = "一份演示 PPT 大纲(6-9 页,每页标题精炼、3-5 个要点,每个要点是完整可讲的一句话,别只写词组)"
    system = (f"你是顶尖行业顾问,正在为一位专业人士交付可直接使用的{what}。主题:「{topic}」。**只基于**下面用户知识库的真实资料。\n"
              f"质量要求(站在使用者角度,像给客户/领导汇报):\n"
              f"①中文,专业、具体、有洞察——每一条都要说清『是什么/为什么重要/接下来怎样』,而不是罗列名词;\n"
              f"②标题和小标题要有信息量、能提炼观点(如『同业存单收量收紧,报价窗口前移』),严禁『概述/总结/背景介绍』这类空话;\n"
              f"③优先用资料里的**真实事实、数字、人名、时间**支撑论点,让它一看就知道是基于真实积累、不是泛泛而谈;\n"
              f"④逻辑要有主线(现状→关键发现→结论/建议),不是零散拼凑;\n"
              f"⑤资料不足处可合理组织框架,但**绝不编造具体数字或事实**,拿不准就不写死。\n"
              f"**只输出 JSON**,不要任何解释或代码围栏,格式:{schema}\n\n===== 用户知识库资料 =====\n{context}")
    # ★产出交付物用 pro 深度模型保质量。关键:max_tokens 给足(16000),否则推理思维链把预算吃光→返空(踩过的坑)。
    out = LLM.chat([{"role": "system", "content": system},
                    {"role": "user", "content": f"主题:{topic}"}],
                   temperature=0.5, max_tokens=16000)
    try:
        return _parse_json(out)
    except Exception:
        # LLM 输出被截断/带杂质 → 让它只重发完整 JSON(更短、别啰嗦)
        fix = LLM.chat([{"role": "system", "content": "你只输出一个完整、合法的 JSON,不要任何解释、不要代码围栏。格式:" + schema + "。控制在合理长度以免被截断。"},
                        {"role": "user", "content": f"主题:{topic}\n\n参考资料:\n{context[:4000]}"}],
                       temperature=0.4, max_tokens=16000)
        return _parse_json(fix)


def _safe(name: str) -> str:
    return "".join(c if c.isalnum() or c in "-_" else "_" for c in name)[:40] or "output"


# ---------- PPT(带设计主题) ----------
def render_ppt(data: dict, tag: str, theme: str = DEFAULT_THEME) -> str:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    th = THEMES.get(theme, THEMES[DEFAULT_THEME])
    C = lambda t: RGBColor(*t)
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9
    prs.slide_height = Emu(6858000)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def paint_bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = C(th["bg"])

    def bar(slide, x, y, w, h, color):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = C(color)
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
        tf.word_wrap = True; tf.vertical_anchor = anchor
        return tf

    def style(p, text, size, color, bold=False, align=PP_ALIGN.LEFT):
        p.text = text; p.alignment = align
        r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = C(color); r.font.name = th["font"]

    # ---- 封面 ----
    s0 = prs.slides.add_slide(blank); paint_bg(s0)
    bar(s0, Emu(838000), Emu(2600000), Emu(1500000), Emu(70000), th["accent"])
    tf = textbox(s0, Emu(838000), Emu(2750000), SW - Emu(1676000), Emu(1600000))
    style(tf.paragraphs[0], data.get("title", "演示"), 40, th["title"], bold=True)
    if data.get("subtitle"):
        style(tf.add_paragraph(), data["subtitle"], 20, th["accent"])
    ftf = textbox(s0, Emu(838000), SH - Emu(700000), SW, Emu(400000))
    style(ftf.paragraphs[0], "由你的第二大脑 · 基于你的知识库生成", 12, th["foot"])

    # ---- 内容页 ----
    for sl in data.get("slides", []):
        s = prs.slides.add_slide(blank); paint_bg(s)
        bar(s, Emu(838000), Emu(560000), Emu(90000), Emu(560000), th["accent"])
        htf = textbox(s, Emu(1050000), Emu(520000), SW - Emu(1888000), Emu(700000))
        style(htf.paragraphs[0], sl.get("title", ""), 26, th["title"], bold=True)
        btf = textbox(s, Emu(1050000), Emu(1500000), SW - Emu(1888000), Emu(4100000))
        bullets = sl.get("bullets", [])
        for j, b in enumerate(bullets):
            p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            style(p, "· " + str(b), 18, th["body"])
            p.space_after = Pt(12)
        if sl.get("note"):
            ntf = textbox(s, Emu(1050000), SH - Emu(760000), SW - Emu(1888000), Emu(500000))
            style(ntf.paragraphs[0], "结论:" + sl["note"], 13, th["accent"], bold=True)

    path = os.path.join(OUT_DIR, f"{_safe(data.get('title', 'ppt'))}_{tag}.pptx")
    prs.save(path)
    return path


# ---------- Word ----------
def render_docx(data: dict, tag: str) -> str:
    from docx import Document
    from docx.shared import Pt, RGBColor
    doc = Document()
    doc.add_heading(data.get("title", "文档"), 0)
    if data.get("subtitle"):
        sp = doc.add_paragraph(); r = sp.add_run(data["subtitle"])
        r.italic = True; r.font.size = Pt(12); r.font.color.rgb = RGBColor(0x66, 0x6c, 0x78)
    for sec in data.get("sections", []):
        doc.add_heading(sec.get("heading", ""), level=1)
        for para in sec.get("paragraphs", []):
            doc.add_paragraph(str(para))
    doc.add_paragraph()
    foot = doc.add_paragraph(); fr = foot.add_run("由你的第二大脑 · 基于你的知识库生成")
    fr.font.size = Pt(9); fr.font.color.rgb = RGBColor(0x99, 0x9c, 0xa6)
    path = os.path.join(OUT_DIR, f"{_safe(data.get('title', 'doc'))}_{tag}.docx")
    doc.save(path)
    return path


# ---------- Excel ----------
def render_xlsx(data: dict, tag: str) -> str:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    wb = Workbook(); ws = wb.active; ws.title = (data.get("title", "Sheet")[:28] or "Sheet")
    cols = data.get("columns", [])
    thin = Side(style="thin", color="D9DCE1")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    if cols:
        ws.append(cols)
        for c in range(1, len(cols) + 1):
            cell = ws.cell(row=1, column=c)
            cell.font = Font(bold=True, color="FFFFFF", size=11)
            cell.fill = PatternFill("solid", fgColor="2563EB")
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = border
    for row in data.get("rows", []):
        ws.append([str(c) for c in row])
    # 自适应列宽 + 边框
    for ci, col in enumerate(ws.columns, 1):
        width = max((len(str(c.value)) for c in col if c.value), default=8)
        ws.column_dimensions[col[0].column_letter].width = min(max(width * 1.6 + 3, 10), 46)
        for c in col:
            if c.row > 1:
                c.border = border
    ws.freeze_panes = "A2"
    path = os.path.join(OUT_DIR, f"{_safe(data.get('title', 'sheet'))}_{tag}.xlsx")
    wb.save(path)
    return path


# ---------- OfficeCLI(Docker)PPT:用主题母版的标题/正文占位符,排版专业 ----------
def render_ppt_officecli(data: dict, tag: str) -> str:
    import subprocess
    fn = f"{_safe(data.get('title', 'ppt'))}_{tag}.pptx"
    path = os.path.join(OUT_DIR, fn)
    work = "/work/" + fn
    OCC = ["occli"]  # wrapper: docker run ... -v OUT_DIR:/work officecli
    subprocess.run(OCC + ["create", work], check=True, capture_output=True, timeout=90)
    cmds = [{"command": "add", "parent": "/", "type": "slide",
             "props": {"layout": "Title Slide", "title": data.get("title", "演示"),
                       "text": data.get("subtitle", "由你的第二大脑生成")}}]
    for sl in data.get("slides", []):
        body = "\n".join(str(b) for b in sl.get("bullets", []))
        if sl.get("note"):
            body += "\n\n结论:" + str(sl["note"])
        cmds.append({"command": "add", "parent": "/", "type": "slide",
                     "props": {"layout": "Title and Content", "title": sl.get("title", ""), "text": body}})
    cf = os.path.join(OUT_DIR, f"_cmds_{tag}.json")
    with open(cf, "w", encoding="utf-8") as f:
        json.dump(cmds, f, ensure_ascii=False)
    try:
        subprocess.run(OCC + ["batch", work, "--input", "/work/" + os.path.basename(cf)],
                       check=True, capture_output=True, timeout=150)
    finally:
        try:
            os.remove(cf)
        except OSError:
            pass
    if not os.path.exists(path):
        raise RuntimeError("officecli 未产出文件")
    _occ_preview(fn)
    return path


OCC = ["occli"]  # wrapper: docker run ... -v OUT_DIR:/work officecli


def _occ_preview(fn):
    """给已产出的文件渲一份 html 预览(可交互,免浏览器,"做完看效果")。"""
    import subprocess
    try:
        base = fn.rsplit(".", 1)[0]
        subprocess.run(OCC + ["view", "/work/" + fn, "html", "-o", "/work/" + base + ".html"],
                       check=True, capture_output=True, timeout=90)
    except Exception:
        pass


def render_docx_officecli(data: dict, tag: str) -> str:
    """Word:OfficeCLI markdown → 样式化 docx。"""
    import subprocess
    fn = f"{_safe(data.get('title', 'doc'))}_{tag}.docx"
    path = os.path.join(OUT_DIR, fn); work = "/work/" + fn
    md = "# " + str(data.get("title", "文档")) + "\n\n"
    if data.get("subtitle"):
        md += "*" + str(data["subtitle"]) + "*\n\n"
    for sec in data.get("sections", []):
        md += "## " + str(sec.get("heading", "")) + "\n\n"
        for para in sec.get("paragraphs", []):
            md += str(para) + "\n\n"
    subprocess.run(OCC + ["create", work], check=True, capture_output=True, timeout=90)
    subprocess.run(OCC + ["add", work, "/body", "--type", "markdown", "--prop", "text=" + md],
                   check=True, capture_output=True, timeout=120)
    if not os.path.exists(path):
        raise RuntimeError("officecli 未产出 docx")
    _occ_preview(fn)
    return path


def render_xlsx_officecli(data: dict, tag: str) -> str:
    """Excel:OfficeCLI import CSV(公式可被求值)。"""
    import subprocess, csv, io
    fn = f"{_safe(data.get('title', 'sheet'))}_{tag}.xlsx"
    path = os.path.join(OUT_DIR, fn); work = "/work/" + fn
    cf = os.path.join(OUT_DIR, f"_data_{tag}.csv")
    with open(cf, "w", encoding="utf-8-sig", newline="") as f:
        w = csv.writer(f)
        cols = data.get("columns", [])
        if cols:
            w.writerow(cols)
        for row in data.get("rows", []):
            w.writerow([str(c) for c in row])
    try:
        subprocess.run(OCC + ["create", work], check=True, capture_output=True, timeout=90)
        subprocess.run(OCC + ["import", work, "/Sheet1", "/work/" + os.path.basename(cf)],
                       check=True, capture_output=True, timeout=120)
    finally:
        try:
            os.remove(cf)
        except OSError:
            pass
    if not os.path.exists(path):
        raise RuntimeError("officecli 未产出 xlsx")
    _occ_preview(fn)
    return path


def generate(LLM, topic: str, sources: list, fmt: str, tag: str, theme: str = DEFAULT_THEME) -> dict:
    data = compose(LLM, topic, sources, fmt)
    if fmt == "word":
        try:
            path = render_docx_officecli(data, tag)
        except Exception:
            path = render_docx(data, tag)
    elif fmt == "excel":
        try:
            path = render_xlsx_officecli(data, tag)
        except Exception:
            path = render_xlsx(data, tag)
    else:
        try:                                  # 优先 OfficeCLI(专业排版),失败退回 python-pptx
            path = render_ppt_officecli(data, tag)
        except Exception:
            path = render_ppt(data, tag, theme)
    fn = os.path.basename(path)
    hn = fn.rsplit(".", 1)[0] + ".html"
    preview = hn if os.path.exists(os.path.join(OUT_DIR, hn)) else None
    return {"file": fn, "title": data.get("title", topic), "format": fmt, "preview": preview}
